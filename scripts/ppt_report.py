# -*- coding: utf-8 -*-
"""Monitoring Visit Analysis - PPT Report Generation."""

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from config import CC, PC, PER_PAGE
from analyze import is_overdue


# PPT color constants
CP = RGBColor(*PC['primary']); CN = RGBColor(*PC['dark']); CA = RGBColor(*PC['danger'])
CG = RGBColor(*PC['success']); CL = RGBColor(*PC['light']); CW = RGBColor(*PC['white'])
CT = RGBColor(*PC['text']); CS = RGBColor(*PC['subtext']); CK = RGBColor(*PC['warning'])
CTL = RGBColor(*PC['teal_light'])

plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'Arial']
plt.rcParams['axes.unicode_minus'] = False


def _savefig(fig, path):
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close(fig)
    return path


def _ashape(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    return s


def _arect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = fill
    return s


def _atxt(slide, l, t, w, h, text, sz=12, bold=False, color=CT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tb


def _akpi(slide, l, t, w, h, num, label, color, sub=None):
    _arect(slide, l, t, w, h, CW)
    _ashape(slide, l + Inches(0.05), t + Inches(0.05), w - Inches(0.1), Inches(0.06), fill=color)
    _atxt(slide, l + Inches(0.15), t + Inches(0.2), w - Inches(0.3), Inches(0.5),
          str(num), sz=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    _atxt(slide, l + Inches(0.1), t + Inches(0.7), w - Inches(0.2), Inches(0.3),
          label, sz=10, color=CS, align=PP_ALIGN.CENTER)
    if sub:
        _atxt(slide, l + Inches(0.1), t + Inches(0.9), w - Inches(0.2), Inches(0.25),
              sub, sz=8, color=CS, align=PP_ALIGN.CENTER)


def _header_bar(slide, title, sw):
    _ashape(slide, Inches(0), Inches(0), sw, Inches(1.0), fill=CN)
    _atxt(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6), title, sz=24, bold=True, color=CW)


def _add_native_table(slide, left, top, width, col_labels, col_widths, cell_data, risk_col):
    """Add a native PPT table with conditional row coloring."""
    nrows = len(cell_data) + 1; ncols = len(col_labels)
    row_h = Inches(0.28); hdr_h = Inches(0.38)
    tshape = slide.shapes.add_table(nrows, ncols, left, top, width, hdr_h + row_h * len(cell_data))
    tbl = tshape.table
    tw = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(width * cw / tw)
    tbl.rows[0].height = hdr_h
    for j, lbl in enumerate(col_labels):
        c = tbl.cell(0, j); c.text = lbl
        c.fill.solid(); c.fill.fore_color.rgb = CP
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = CW; p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    for i, rd in enumerate(cell_data):
        tbl.rows[i + 1].height = row_h
        risk = rd[risk_col]
        for j, val in enumerate(rd):
            c = tbl.cell(i + 1, j); c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = (RGBColor(*PC['hr_bg']) if risk == '\u9ad8\u98ce\u9669' else
                                     RGBColor(*PC['at_bg']) if risk == '\u5173\u6ce8' else
                                     RGBColor(*PC['alt_row']) if i % 2 == 1 else CW)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(9); p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
            if j == risk_col and risk == '\u9ad8\u98ce\u9669':
                p.font.color.rgb = RGBColor(*PC['hr_text']); p.font.bold = True
            elif j == risk_col and risk == '\u5173\u6ce8':
                p.font.color.rgb = RGBColor(*PC['at_text']); p.font.bold = True
            else:
                p.font.color.rgb = CT
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.margin_left = Inches(0.04); c.margin_right = Inches(0.04)
    return hdr_h + row_h * len(cell_data)


def _generate_charts(month_stats, work_dir):
    """Generate all matplotlib chart images. Returns dict of paths."""
    months_sorted = sorted(month_stats.keys())
    ms_short = [m.replace('2025-', '25/').replace('2026-', '26/').replace('2024-', '24/') for m in months_sorted]
    x = np.arange(len(ms_short))
    paths = {}

    # Chart 1: Monthly visits bar chart
    fig, ax = plt.subplots(figsize=(12, 4.5))
    v1 = [month_stats[m]['total'] for m in months_sorted]
    v2 = [month_stats[m]['submitted'] for m in months_sorted]
    v3 = [month_stats[m]['finalized'] for m in months_sorted]
    bw = 0.25
    b1 = ax.bar(x - bw, v1, bw, label='\u8bbf\u89c6\u6b21\u6570', color=CC['C1'], zorder=3)
    b2 = ax.bar(x, v2, bw, label='\u5df2\u9012\u4ea4', color=CC['C2'], zorder=3)
    b3 = ax.bar(x + bw, v3, bw, label='\u5df2\u5b9a\u7a3f', color=CC['C3'], zorder=3)
    for bars in [b1, b2, b3]:
        for bar in bars:
            h = bar.get_height()
            if h > 0: ax.text(bar.get_x() + bar.get_width() / 2., h + 0.3, str(int(h)),
                              ha='center', va='bottom', fontsize=8, fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(ms_short, fontsize=11)
    ax.set_ylabel('\u6b21\u6570', fontsize=12); ax.set_ylim(0, max(v1) * 1.3 if v1 else 10)
    ax.legend(fontsize=11, loc='upper left', ncol=3)
    ax.set_title('\u6708\u5ea6\u8bbf\u89c6 / \u9012\u4ea4 / \u5b9a\u7a3f\u6b21\u6570', fontsize=15, fontweight='bold', pad=12)
    ax.grid(axis='y', alpha=0.3); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    paths['monthly'] = _savefig(fig, os.path.join(work_dir, 'monthly_visits.png'))

    # Chart 2: Overdue trend (counts + rates)
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(12, 5.5), gridspec_kw={'height_ratios': [1, 1]})
    sub_od = [month_stats[m]['sub_od'] for m in months_sorted]
    fin_od = [month_stats[m]['fin_od'] for m in months_sorted]
    fu_od = [month_stats[m]['fu_od'] for m in months_sorted]
    bw2 = 0.25
    ax1.bar(x - bw2, sub_od, bw2, label='\u9012\u4ea4\u8d85\u671f', color=CC['C1'], zorder=3)
    ax1.bar(x, fin_od, bw2, label='\u5b9a\u7a3f\u8d85\u671f', color=CC['C4'], zorder=3)
    ax1.bar(x + bw2, fu_od, bw2, label='\u8ddf\u8fdb\u51fd\u8d85\u671f', color=CC['C3'], zorder=3)
    for i in range(len(x)):
        for val, offset in [(sub_od[i], -bw2), (fin_od[i], 0), (fu_od[i], bw2)]:
            if val > 0: ax1.text(x[i] + offset, val + 0.1, str(val), ha='center', fontsize=9, fontweight='bold')
    ax1.set_xticks(x); ax1.set_xticklabels(ms_short, fontsize=10)
    ax1.set_ylabel('\u8d85\u671f\u6b21\u6570', fontsize=11)
    ax1.set_ylim(0, max(max(sub_od + fin_od + fu_od) + 2, 5))
    ax1.legend(fontsize=10, ncol=3, loc='upper right')
    ax1.set_title('\u6708\u5ea6\u8d85\u671f\u6b21\u6570', fontsize=13, fontweight='bold', pad=8)
    ax1.grid(axis='y', alpha=0.3); ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)
    sub_rate = [month_stats[m]['sub_od'] / month_stats[m]['submitted'] * 100 if month_stats[m]['submitted'] else 0 for m in months_sorted]
    fin_rate = [month_stats[m]['fin_od'] / month_stats[m]['finalized'] * 100 if month_stats[m]['finalized'] else 0 for m in months_sorted]
    fu_rate = [month_stats[m]['fu_od'] / month_stats[m]['finalized'] * 100 if month_stats[m]['finalized'] else 0 for m in months_sorted]
    ax2.plot(x, sub_rate, 'o-', color=CC['C1'], linewidth=2, markersize=6, label='\u9012\u4ea4\u8d85\u671f\u7387')
    ax2.plot(x, fin_rate, 's-', color=CC['C4'], linewidth=2, markersize=6, label='\u5b9a\u7a3f\u8d85\u671f\u7387')
    ax2.plot(x, fu_rate, '^-', color=CC['C3'], linewidth=2, markersize=6, label='\u8ddf\u8fdb\u51fd\u8d85\u671f\u7387')
    for i in range(len(x)):
        for val, c in [(sub_rate[i], CC['C1']), (fin_rate[i], CC['C4']), (fu_rate[i], CC['C3'])]:
            if val > 0: ax2.text(x[i], val + 2, f'{val:.0f}%', ha='center', fontsize=8, fontweight='bold', color=c)
    ax2.set_xticks(x); ax2.set_xticklabels(ms_short, fontsize=10)
    ax2.set_ylabel('\u8d85\u671f\u7387 %', fontsize=11)
    ax2.set_ylim(-5, max(max(sub_rate + fin_rate + fu_rate) + 15, 20))
    ax2.legend(fontsize=10, ncol=3, loc='upper right')
    ax2.set_title('\u6708\u5ea6\u8d85\u671f\u7387\u8d8b\u52bf', fontsize=13, fontweight='bold', pad=8)
    ax2.grid(axis='y', alpha=0.3); ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)
    plt.tight_layout()
    paths['overdue_trend'] = _savefig(fig, os.path.join(work_dir, 'overdue_trend.png'))

    # Chart 3: Avg WD trend
    fig, ax = plt.subplots(figsize=(12, 3.0))
    avg_sub = [sum(month_stats[m]['sub_wd']) / len(month_stats[m]['sub_wd']) if month_stats[m]['sub_wd'] else 0 for m in months_sorted]
    avg_fin = [sum(month_stats[m]['fin_wd']) / len(month_stats[m]['fin_wd']) if month_stats[m]['fin_wd'] else 0 for m in months_sorted]
    ax.plot(x, avg_sub, 'o-', color=CC['C1'], linewidth=2.5, markersize=7, label='\u5e73\u5747\u9012\u4ea4WD', zorder=3)
    ax.plot(x, avg_fin, 's-', color=CC['C4'], linewidth=2.5, markersize=7, label='\u5e73\u5747\u5b9a\u7a3fWD', zorder=3)
    ax.axhline(y=5, color=CC['C1'], linestyle='--', alpha=0.4)
    ax.axhline(y=10, color=CC['C4'], linestyle='--', alpha=0.4)
    ax.text(len(x) - 0.5, 5.2, '\u9012\u4ea4\u6807\u51c65WD', fontsize=8, color=CC['C1'], alpha=0.6)
    ax.text(len(x) - 0.5, 10.2, '\u5b9a\u7a3f\u6807\u516810WD', fontsize=8, color=CC['C4'], alpha=0.6)
    ax.set_xticks(x); ax.set_xticklabels(ms_short, fontsize=9)
    ax.set_ylabel('\u5de5\u4f5c\u65e5', fontsize=10); ax.set_ylim(0, 15)
    ax.legend(fontsize=9, loc='upper left')
    ax.set_title('\u5e73\u5747\u5de5\u4f5c\u65e5\u8d8b\u52bf\uff08\u6708\u5ea6\uff09', fontsize=13, fontweight='bold', pad=10)
    ax.grid(axis='y', alpha=0.3); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    paths['trend_wd'] = _savefig(fig, os.path.join(work_dir, 'trend_wd.png'))

    # Chart 4: Archive pie
    total_archived = sum(s['archived'] for s in month_stats.values())
    total_visited = sum(s['total'] for s in month_stats.values())
    fig, ax = plt.subplots(figsize=(3.5, 3.0))
    ax.pie([total_archived, total_visited - total_archived],
           labels=[f'\u5df2\u5f52\u6863 ({total_archived})', f'\u672a\u5f52\u6863 ({total_visited - total_archived})'],
           colors=[CC['C5'], '#E0E0E0'], autopct='%1.0f%%', startangle=90, textprops={'fontsize': 11})
    ax.set_title('\u62a5\u544a\u5f52\u6863\u72b6\u6001', fontsize=13, fontweight='bold', pad=8)
    paths['pie'] = _savefig(fig, os.path.join(work_dir, 'pie_archive.png'))

    # Chart 5: Archive rate trend
    fig, ax = plt.subplots(figsize=(12, 2.0))
    arch_rates = [month_stats[m]['archived'] / month_stats[m]['total'] * 100 for m in months_sorted]
    ax.fill_between(x, arch_rates, alpha=0.3, color=CC['C5'])
    ax.plot(x, arch_rates, 'o-', color=CC['C5'], linewidth=2.5, markersize=6)
    for i, v in enumerate(arch_rates):
        ax.text(i, v + 2, f'{v:.0f}%', ha='center', fontsize=9, fontweight='bold', color=CC['C5'])
    ax.set_xticks(x); ax.set_xticklabels(ms_short, fontsize=10)
    ax.set_ylim(0, 115); ax.axhline(y=100, color='#999', linestyle='--', alpha=0.3)
    ax.set_title('\u6708\u5ea6\u5f52\u6863\u7387\u8d8b\u52bf', fontsize=13, fontweight='bold', pad=8)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False); ax.grid(axis='y', alpha=0.2)
    paths['arch_rate'] = _savefig(fig, os.path.join(work_dir, 'archive_rate.png'))

    return paths


def generate_ppt(visited, cra_stats, month_stats, site_stats, summary, work_dir, output_path, today,
                 project_name='Monitoring Visit'):
    """Generate the full PPT visualization report."""
    sm = summary
    charts = _generate_charts(month_stats, work_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW = Inches(13.333); SH = Inches(7.5)

    # ---- Slide 1: COVER ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CN)
    _ashape(s, Inches(0), Inches(0), Inches(0.15), SH, fill=CP)
    _ashape(s, Inches(0), Inches(4.8), SW, Inches(0.04), fill=CP)
    _atxt(s, Inches(1.2), Inches(1.5), Inches(10), Inches(1), 'MONITORING VISIT', sz=18, color=CTL)
    _atxt(s, Inches(1.2), Inches(2.1), Inches(10), Inches(1.2),
          '\u76d1\u67e5\u8bbf\u89c6\u8ffd\u8e2a\u5206\u6790\u62a5\u544a', sz=40, bold=True, color=CW)
    _atxt(s, Inches(1.2), Inches(3.5), Inches(10), Inches(0.5), project_name, sz=16, color=CTL)
    _atxt(s, Inches(1.2), Inches(5.2), Inches(8), Inches(0.4),
          f'\u6570\u636e\u622a\u6b62\u65e5\u671f\uff1a{today}', sz=14, color=CS)
    _atxt(s, Inches(1.2), Inches(5.7), Inches(10), Inches(0.7),
          '\u8ddf\u8fdb\u51fd\u8d85\u671f = min(\u8bbf\u89c6\u7ed3\u675f+10WD, \u5b9a\u7a3f+1WD)\uff0c\u4ec5\u5df2\u5b9a\u7a3f\u62a5\u544a\n'
          '\u6f0f\u767b\u8bb0 = \u5df2\u5b9a\u7a3f + \u8ddf\u8fdb\u51fd\u622a\u6b62\u5df2\u8fc7 + \u672a\u53d1\u9001 | \u98ce\u9669 = \u4efb\u4e00\u7ef4\u5ea6\u22652\u2192\u9ad8\u98ce\u9669',
          sz=11, color=CS)

    # ---- Slide 2: DASHBOARD ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
    _header_bar(s, '\u9879\u76ee\u6982\u89c8 DASHBOARD', SW)
    y1 = Inches(1.3); cw = Inches(2.4); ch = Inches(1.2); gap = Inches(0.3); x0 = Inches(0.6)
    _akpi(s, x0, y1, cw, ch, str(sm['total_records']), '\u8bbf\u89c6\u8bb0\u5f55\u603b\u6570', CP)
    _akpi(s, x0 + cw + gap, y1, cw, ch, str(sm['total_visited']), '\u5df2\u5b8c\u6210\u8bbf\u89c6', CP,
          f"\u5b8c\u6210\u7387 {sm['total_visited']/sm['total_records']*100:.1f}%" if sm['total_records'] else '')
    _akpi(s, x0 + (cw + gap) * 2, y1, cw, ch, str(sm['total_sub_od']), '\u9012\u4ea4\u8d85\u671f',
          CG if sm['total_sub_od'] == 0 else CA, '\u96f6\u8d85\u671f' if sm['total_sub_od'] == 0 else '')
    _akpi(s, x0 + (cw + gap) * 3, y1, cw, ch, str(sm['total_fin_od']), '\u5b9a\u7a3f\u8d85\u671f',
          CA if sm['total_fin_od'] > 0 else CG)
    _akpi(s, x0 + (cw + gap) * 4 - Inches(0.05), y1, cw + Inches(0.1), ch,
          str(sm['total_fu_od']), '\u8ddf\u8fdb\u51fd\u8d85\u671f', CA if sm['total_fu_od'] > 0 else CG)
    y2 = Inches(2.8)
    _akpi(s, x0, y2, cw, ch, str(sm['total_missing']), '\u7591\u4f3c\u6f0f\u767b\u8bb0',
          CK if sm['total_missing'] > 0 else CG, '\u9700\u6838\u5b9e' if sm['total_missing'] > 0 else '')
    _akpi(s, x0 + cw + gap, y2, cw, ch, str(sm['total_archived']), '\u5df2\u5f52\u6863', CG,
          f"\u5f52\u6863\u7387 {sm['total_archived']/sm['total_visited']*100:.1f}%" if sm['total_visited'] else '')
    _akpi(s, x0 + (cw + gap) * 2, y2, cw, ch, str(len(sm['high_risk_cras'])), '\u9ad8\u98ce\u9669CRA',
          CA if sm['high_risk_cras'] else CG,
          ', '.join(sm['high_risk_cras']) if sm['high_risk_cras'] else '\u65e0')
    _akpi(s, x0 + (cw + gap) * 3, y2, cw, ch, str(len(sm['high_risk_sites'])), '\u9ad8\u98ce\u9669\u7814\u7a76\u4e2d\u5fc3',
          CA if sm['high_risk_sites'] else CG)
    _akpi(s, x0 + (cw + gap) * 4 - Inches(0.05), y2, cw + Inches(0.1), ch,
          f"{sm['unique_cras']} / {sm['unique_sites']}", 'CRA / \u7814\u7a76\u4e2d\u5fc3', CTL)
    s.shapes.add_picture(charts['pie'], Inches(0.6), Inches(4.5), Inches(4.0), Inches(2.8))
    s.shapes.add_picture(charts['trend_wd'], Inches(5.0), Inches(4.6), Inches(7.8), Inches(2.5))

    # ---- Slide 3: MONTHLY TREND ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
    _header_bar(s, '\u6708\u5ea6\u8bbf\u89c6\u8d8b\u52bf MONTHLY TREND', SW)
    s.shapes.add_picture(charts['monthly'], Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.0))
    s.shapes.add_picture(charts['arch_rate'], Inches(0.4), Inches(5.3), Inches(12.5), Inches(2.0))

    # ---- Slide 4: OVERDUE TREND ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
    _header_bar(s, '\u6708\u5ea6\u8d85\u671f\u8d8b\u52bf OVERDUE TREND', SW)
    s.shapes.add_picture(charts['overdue_trend'], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5))
    _arect(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.5), CW)
    _atxt(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
          '\u8bf4\u660e\uff1a\u8ddf\u8fdb\u51fd\u8d85\u671f\u7edf\u8ba1\u4ec5\u9488\u5bf9\u5df2\u5b9a\u7a3f\u7684\u76d1\u67e5\u62a5\u544a\uff0c\u57fa\u4e8emin(\u8bbf\u89c6\u7ed3\u675f+10WD, \u5b9a\u7a3f+1WD)\u89c4\u5219\u3002',
          sz=10, color=CT)

    # ---- Slide 5+: CRA TABLES ----
    cra_sorted = sorted(cra_stats.items(), key=lambda x: x[1]['total'], reverse=True)
    cra_pages = [cra_sorted[i:i + PER_PAGE] for i in range(0, len(cra_sorted), PER_PAGE)]
    cra_od_notes = []
    for name, sd in cra_sorted:
        for d in sd['od_details']:
            cra_od_notes.append(f"{name}: {', '.join(d['items'])} (\u4e2d\u5fc3{d['site']}, {d['visit']})")

    cra_col_labels = ['CRA', '\u8bbf\u89c6', '\u9012\u4ea4\u8d85\u671f', '\u5b9a\u7a3f\u8d85\u671f', '\u8ddf\u8fdb\u51fd\n\u8d85\u671f', '\u6f0f\u767b\u8bb0',
                      '\u5e73\u5747\n\u9012\u4ea4WD', '\u5e73\u5747\n\u5b9a\u7a3fWD', '\u5df2\u5f52\u6863', '\u98ce\u9669\u7b49\u7ea7']
    cra_col_w = [2.5, 0.8, 0.95, 0.95, 0.95, 0.8, 1.05, 1.05, 0.8, 1.05]
    for pi, page in enumerate(cra_pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
        _header_bar(s, f'CRA \u7ee9\u6548\u5206\u6790 ({pi + 1}/{len(cra_pages)})', SW)
        cd = []
        for cname, st in page:
            avg_s = f"{sum(st['sub_wd'])/len(st['sub_wd']):.1f}" if st['sub_wd'] else '-'
            avg_f = f"{sum(st['fin_wd'])/len(st['fin_wd']):.1f}" if st['fin_wd'] else '-'
            cd.append([cname, str(st['total']), str(st['sub_od']), str(st['fin_od']),
                       str(st['fu_od']), str(st['missing']), avg_s, avg_f, str(st['archived']), st['risk']])
        tbl_h = _add_native_table(s, Inches(0.3), Inches(1.15), Inches(12.7), cra_col_labels, cra_col_w, cd, 9)
        note_y = Inches(1.15) + tbl_h + Inches(0.15)
        _arect(s, Inches(0.3), note_y, Inches(12.7), SH - note_y - Inches(0.15), CW)
        _ashape(s, Inches(0.3), note_y, Inches(0.08), SH - note_y - Inches(0.15), fill=CA)
        _atxt(s, Inches(0.6), note_y + Inches(0.05), Inches(3), Inches(0.3),
              '\u8d85\u671f/\u6f0f\u767b\u8bb0\u8be6\u60c5', sz=11, bold=True, color=CA)
        page_cras = {cn for cn, _ in page}
        pn = [n for n in cra_od_notes if n.split(':')[0] in page_cras]
        nt = ' | '.join(pn[:8]) if pn else '\u8be5\u9875CRA\u65e0\u8d85\u671f/\u6f0f\u767b\u8bb0\u8bb0\u5f55'
        _atxt(s, Inches(0.6), note_y + Inches(0.3), Inches(12.2), SH - note_y - Inches(0.5), nt, sz=9, color=CS)

    # ---- Site TABLES ----
    site_sorted = sorted(site_stats.items(), key=lambda x: x[1]['total'], reverse=True)
    site_pages = [site_sorted[i:i + PER_PAGE] for i in range(0, len(site_sorted), PER_PAGE)]
    site_od_notes = []
    for sid, sd in site_sorted:
        for d in sd['od_details']:
            site_od_notes.append(f"\u4e2d\u5fc3{sid}: {d['cra']} {', '.join(d['items'])} ({d['visit']})")

    site_col_labels = ['\u7f16\u53f7', '\u4e2d\u5fc3\u7b80\u79f0', '\u8bbf\u89c6', '\u9012\u4ea4\u8d85\u671f', '\u5b9a\u7a3f\u8d85\u671f', '\u8ddf\u8fdb\u51fd\n\u8d85\u671f',
                       '\u6f0f\u767b\u8bb0', '\u5e73\u5747\n\u9012\u4ea4WD', '\u5e73\u5747\n\u5b9a\u7a3fWD', '\u5df2\u5f52\u6863', '\u98ce\u9669\u7b49\u7ea7']
    site_col_w = [1.0, 2.2, 0.7, 0.9, 0.9, 0.9, 0.7, 1.0, 1.0, 0.7, 1.0]
    for pi, page in enumerate(site_pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
        _header_bar(s, f'\u7814\u7a76\u4e2d\u5fc3\u5206\u6790 ({pi + 1}/{len(site_pages)})', SW)
        cd = []
        for sid_v, st in page:
            sn = st['name'][:8] + ('...' if len(st['name']) > 8 else '')
            avg_s = f"{sum(st['sub_wd'])/len(st['sub_wd']):.1f}" if st['sub_wd'] else '-'
            avg_f = f"{sum(st['fin_wd'])/len(st['fin_wd']):.1f}" if st['fin_wd'] else '-'
            cd.append([sid_v, sn, str(st['total']), str(st['sub_od']), str(st['fin_od']),
                       str(st['fu_od']), str(st['missing']), avg_s, avg_f, str(st['archived']), st['risk']])
        tbl_h = _add_native_table(s, Inches(0.3), Inches(1.15), Inches(12.7), site_col_labels, site_col_w, cd, 10)
        note_y = Inches(1.15) + tbl_h + Inches(0.15)
        _arect(s, Inches(0.3), note_y, Inches(12.7), SH - note_y - Inches(0.15), CW)
        _ashape(s, Inches(0.3), note_y, Inches(0.08), SH - note_y - Inches(0.15), fill=CK)
        _atxt(s, Inches(0.6), note_y + Inches(0.05), Inches(3), Inches(0.3),
              '\u8d85\u671f/\u6f0f\u767b\u8bb0\u8be6\u60c5', sz=11, bold=True, color=CK)
        page_sids = {sv for sv, _ in page}
        pn = [n for n in site_od_notes if any(f'\u4e2d\u5fc3{sv}' in n for sv in page_sids)]
        nt = ' | '.join(pn[:8]) if pn else '\u8be5\u9875\u7814\u7a76\u4e2d\u5fc3\u65e0\u8d85\u671f/\u6f0f\u767b\u8bb0\u8bb0\u5f55'
        _atxt(s, Inches(0.6), note_y + Inches(0.3), Inches(12.2), SH - note_y - Inches(0.5), nt, sz=9, color=CS)

    # ---- ISSUES DETAIL ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CL)
    _header_bar(s, '\u8d85\u671f\u8bb0\u5f55 & \u6f0f\u767b\u8bb0\u660e\u7ec6 ISSUES', SW)
    fin_od_recs = [r for r in visited if is_overdue(r.get('final_overdue'))]
    _arect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.5), CW)
    _ashape(s, Inches(0.5), Inches(1.3), Inches(0.1), Inches(2.5), fill=CA)
    _atxt(s, Inches(0.9), Inches(1.4), Inches(5), Inches(0.35),
          f'\u5b9a\u7a3f\u8d85\u671f\u8bb0\u5f55\uff08{sm["total_fin_od"]}\u6761\uff09', sz=15, bold=True, color=CA)
    yo = Inches(1.85)
    for r in fin_od_recs[:4]:
        _atxt(s, Inches(0.9), yo, Inches(5.4), Inches(0.25),
              f"\u4e2d\u5fc3{r['site_no']} | {r['site_name'][:10]} | CRA: {r['cra']}", sz=11, bold=True, color=CT)
        yo += Inches(0.28)
        _atxt(s, Inches(0.9), yo, Inches(5.4), Inches(0.25),
              f"\u8bbf\u89c6: {r['visit_end']} | \u5b9a\u7a3fWD: {r['final_wd']}\u5929 | \u8d85\u671f{r['final_wd']-10}WD", sz=10, color=CS)
        yo += Inches(0.35)
    missing_recs = [r for r in visited if r.get('missing_reg')]
    _arect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.5), CW)
    _ashape(s, Inches(6.8), Inches(1.3), Inches(0.1), Inches(2.5), fill=CK)
    _atxt(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.35),
          f'\u7591\u4f3c\u6f0f\u767b\u8bb0\uff08{sm["total_missing"]}\u6761\uff09', sz=15, bold=True, color=CK)
    yo2 = Inches(1.85)
    for r in missing_recs[:4]:
        _atxt(s, Inches(7.2), yo2, Inches(5.4), Inches(0.25),
              f"\u4e2d\u5fc3{r['site_no']} | {r['site_name'][:10]} | CRA: {r['cra']}", sz=11, bold=True, color=CT)
        yo2 += Inches(0.28)
        fu_st = r['followup_status'] if r['followup_status'] not in ('', 'None') else '\u7a7a'
        _atxt(s, Inches(7.2), yo2, Inches(5.4), Inches(0.25),
              f"\u8bbf\u89c6: {r['visit_end']} | \u5df2\u5b9a\u7a3f | \u8ddf\u8fdb\u51fd: {fu_st} | \u622a\u6b62: {r['followup_deadline']}", sz=10, color=CS)
        yo2 += Inches(0.35)
    fu_od_recs = [r for r in visited if is_overdue(r.get('followup_overdue'))]
    _arect(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.6), CW)
    _ashape(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.06), fill=CP)
    _atxt(s, Inches(0.8), Inches(4.25), Inches(5), Inches(0.35),
          f'\u8ddf\u8fdb\u51fd\u8d85\u671f\uff08{sm["total_fu_od"]}\u6761\uff09', sz=15, bold=True, color=CP)
    if fu_od_recs:
        yo3 = Inches(4.65)
        for r in fu_od_recs[:3]:
            _atxt(s, Inches(0.8), yo3, Inches(12), Inches(0.25),
                  f"\u4e2d\u5fc3{r['site_no']} {r['site_name'][:10]} | CRA: {r['cra']} | \u8bbf\u89c6: {r['visit_end']} | "
                  f"\u622a\u6b62: {r['followup_deadline']} | \u8d85\u671f{r.get('followup_wd', 0)}WD", sz=10, color=CT)
            yo3 += Inches(0.3)
    else:
        _atxt(s, Inches(0.8), Inches(4.65), Inches(12), Inches(0.25),
              '\u65e0\u8ddf\u8fdb\u51fd\u8d85\u671f\u8bb0\u5f55', sz=11, color=CG, bold=True)
    # Rules
    _arect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.4), CW)
    _ashape(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.06), fill=CP)
    _atxt(s, Inches(0.8), Inches(6.05), Inches(3), Inches(0.3), '\u8d85\u671f\u5224\u5b9a\u89c4\u5219', sz=13, bold=True, color=CP)
    rules = [('\u9012\u4ea4\u8d85\u671f', '\u8bbf\u89c6\u7ed3\u675f + 5WD'), ('\u5b9a\u7a3f\u8d85\u671f', '\u8bbf\u89c6\u7ed3\u675f + 10WD'),
             ('\u8ddf\u8fdb\u51fd\u8d85\u671f', 'min(\u8bbf\u89c6\u7ed3\u675f+10WD, \u5b9a\u7a3f\u65e5\u671f+1WD)\uff0c\u4ec5\u5df2\u5b9a\u7a3f\u62a5\u544a'),
             ('\u6f0f\u767b\u8bb0', '\u5df2\u5b9a\u7a3f + \u8ddf\u8fdb\u51fd\u622a\u6b62\u5df2\u8fc7 + \u8ddf\u8fdb\u51fd\u672a\u53d1\u9001'),
             ('\u98ce\u9669\u7b49\u7ea7', '\u4efb\u4e00\u7ef4\u5ea6(\u9012\u4ea4/\u5b9a\u7a3f/\u8ddf\u8fdb\u51fd\u8d85\u671f,\u6f0f\u767b\u8bb0) \u22652\u2192\u9ad8\u98ce\u9669, =1\u2192\u5173\u6ce8')]
    yr = Inches(6.35)
    for label, desc in rules:
        _ashape(s, Inches(1.0), yr + Inches(0.06), Inches(0.1), Inches(0.1), fill=CP)
        _atxt(s, Inches(1.3), yr, Inches(1.8), Inches(0.22), label, sz=10, bold=True, color=CT)
        _atxt(s, Inches(3.1), yr, Inches(9.5), Inches(0.22), desc, sz=10, color=CS)
        yr += Inches(0.22)

    # ---- SUMMARY ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ashape(s, Inches(0), Inches(0), SW, SH, fill=CN)
    _ashape(s, Inches(0), Inches(0), Inches(0.15), SH, fill=CP)
    _atxt(s, Inches(1.2), Inches(0.8), Inches(10), Inches(0.5), 'SUMMARY', sz=16, color=CTL)
    _atxt(s, Inches(1.2), Inches(1.2), Inches(10), Inches(0.8), '\u603b\u7ed3\u4e0e\u5efa\u8bae', sz=36, bold=True, color=CW)
    findings = [
        ('01', f'\u62a5\u544a\u9012\u4ea4\u8d85\u671f {sm["total_sub_od"]} \u6b21',
         '\u96f6\u8d85\u671f\uff0c\u5e73\u5747\u9012\u4ea4\u5de5\u4f5c\u65e5\u4f18\u4e8e5WD\u6807\u51c6\u3002' if sm['total_sub_od'] == 0 else f'\u5171 {sm["total_sub_od"]} \u6b21\u8d85\u671f\uff0c\u8be6\u89c1\u8d85\u671f\u660e\u7ec6\u3002',
         CG if sm['total_sub_od'] == 0 else CA),
        ('02', f'\u5b9a\u7a3f\u8d85\u671f {sm["total_fin_od"]} \u6b21',
         f'\u5171 {sm["total_fin_od"]} \u6b21\u5b9a\u7a3f\u8d85\u671f\u3002' if sm['total_fin_od'] else '\u5168\u90e8\u62a5\u544a\u5747\u5728\u622a\u6b62\u65e5\u5185\u5b8c\u6210\u5b9a\u7a3f\u3002',
         CA if sm['total_fin_od'] else CG),
        ('03', f'\u8ddf\u8fdb\u51fd\u8d85\u671f {sm["total_fu_od"]} \u6b21',
         '\u6240\u6709\u5df2\u5b9a\u7a3f\u62a5\u544a\u8ddf\u8fdb\u51fd\u6309\u65f6\u53d1\u9001\u3002' if sm['total_fu_od'] == 0 else '\u8be6\u89c1\u8d85\u671f\u660e\u7ec6\u3002',
         CG if sm['total_fu_od'] == 0 else CA),
        ('04', f'{sm["total_missing"]} \u6761\u7591\u4f3c\u6f0f\u767b\u8bb0\u9700\u6838\u5b9e',
         '\u6d89\u53ca ' + ', '.join(set(r['cra'] for r in visited if r.get('missing_reg'))) + '\uff0c\u5efa\u8bae\u786e\u8ba4\u8ddf\u8fdb\u51fd\u767b\u8bb0\u60c5\u51b5\u3002' if sm['total_missing'] else '\u65e0\u6f0f\u767b\u8bb0\u3002',
         CK if sm['total_missing'] else CG),
        ('05', '\u9ad8\u98ce\u9669CRA/\u7814\u7a76\u4e2d\u5fc3',
         '\u9ad8\u98ce\u9669CRA: ' + (", ".join(sm["high_risk_cras"]) if sm["high_risk_cras"] else "\u65e0") + ' | '
         '\u9ad8\u98ce\u9669\u4e2d\u5fc3: ' + (", ".join(sm["high_risk_sites"]) if sm["high_risk_sites"] else "\u65e0"),
         CA if (sm['high_risk_cras'] or sm['high_risk_sites']) else CG),
    ]
    yf = Inches(2.3)
    for num, title, desc, color in findings:
        _arect(s, Inches(1.2), yf, Inches(11), Inches(0.95), RGBColor(*PC['dark_card']))
        _ashape(s, Inches(1.2), yf, Inches(0.08), Inches(0.95), fill=color)
        _atxt(s, Inches(1.6), yf + Inches(0.08), Inches(0.5), Inches(0.35), num, sz=20, bold=True, color=color)
        _atxt(s, Inches(2.3), yf + Inches(0.08), Inches(9.5), Inches(0.35), title, sz=16, bold=True, color=CW)
        _atxt(s, Inches(2.3), yf + Inches(0.48), Inches(9.5), Inches(0.42), desc, sz=12, color=CS)
        yf += Inches(1.05)
    _ashape(s, Inches(0), Inches(6.9), SW, Inches(0.04), fill=CP)
    _atxt(s, Inches(1.2), Inches(7.0), Inches(10), Inches(0.35),
          f'{project_name} | \u6570\u636e\u622a\u6b62 {today}', sz=10, color=CS)

    prs.save(output_path)
    return output_path
